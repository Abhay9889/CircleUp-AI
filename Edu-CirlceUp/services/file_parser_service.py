import logging
import pathlib


logger =logging.getLogger(__name__)
def extract_text(file_bytes:bytes,file_type:str)->str:
    ext=file_type.lower().strip(".")
    if ext=="txt":
        return parse_txt(file_bytes)
    elif ext=="pdf":
        return parse_pdf(file_bytes)
    elif ext=="docx":
        return parse_docx(file_bytes)    
    elif ext=="pptx":
        return parse_pptx(file_bytes)      
    else:
        logger.warning("Unsupported file type: %s",ext)
        return ""

def parse_txt(data:bytes)->str:
    try:
        return data.decode("utf-8",errors="ignore")
    except Exception as e:
        logger.error("TXT parse error %s",e)
        return ""

def parse_pdf(data: bytes) -> str:
    try:
        import io
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        pages = [page.extract_text() or "" for page in reader.pages]

        return "\n".join(pages)

    except Exception as e:
        logger.error("PDF parse error: %s", e)
        return ""
def parse_docx(data:bytes)->str:
    try:
        import io
        from docx import Document
        doc  = Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        logger.error("DOCX parse error: %s", e)
        return ""
    
def parse_pptx(data:bytes)->str:
    try:
        import io
        from pptx import Presentation
        prs   = Presentation(io.BytesIO(data))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n".join(texts)
    except Exception as e:
        logger.error("PPTX parse error: %s", e)
        return ""
